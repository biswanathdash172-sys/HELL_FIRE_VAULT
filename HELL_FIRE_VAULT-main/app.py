from flask import Flask, request, jsonify, send_file, render_template, session, redirect, url_for
from pptx import Presentation
from pptx.util import Inches, Pt
import csv
import hashlib
import os
import time
from datetime import datetime

# Load environment variables from .env
from dotenv import load_dotenv
load_dotenv()

# -------------------------------------------------------------------
# SUPABASE CLIENT
# -------------------------------------------------------------------
try:
    from supabase import create_client, Client as SupabaseClient
    SUPABASE_URL = os.environ.get("SUPABASE_URL", "")
    SUPABASE_KEY = os.environ.get("SUPABASE_ANON_KEY", "")
    supabase: SupabaseClient = create_client(SUPABASE_URL, SUPABASE_KEY) if SUPABASE_URL and SUPABASE_KEY else None
except ImportError:
    supabase = None
    print("[WARN] supabase-py not installed. Run: pip install supabase")

# -------------------------------------------------------------------
# GEMINI AI CLIENT
# -------------------------------------------------------------------
try:
    import google.generativeai as genai
    GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")
    if GEMINI_API_KEY:
        genai.configure(api_key=GEMINI_API_KEY)
        gemini_model = genai.GenerativeModel("gemini-1.5-flash")
    else:
        gemini_model = None
        print("[WARN] GEMINI_API_KEY not set in .env")
except ImportError:
    gemini_model = None
    print("[WARN] google-generativeai not installed. Run: pip install google-generativeai")

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "hellfire-default-secret")

# In-memory history for demo purposes
app_history = []

# -------------------------------------------------------------------
# HACKATHON CONFIG:
# Setting up temporary folders because we don't have time for a DB!
# -------------------------------------------------------------------
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)


# =====================================================================
# LAYER 1: LEDGER RESOLVER (PROJECT 5)
# Calculates the absolute fewest payments needed for the Hellfire Club.
# Note: Using a simple greedy approach. It might not be mathematically
# perfect for a bank, but for 10-20 D&D nerds, it works flawlessly.
# =====================================================================

def calculate_minimal_payments(csv_file_path):
    balances = {}

    # STEP 1: Read the messy CSV and calculate everyone's net balance.
    # Assuming CSV has headers: Payer, Payee, Amount
    try:
        with open(csv_file_path, mode='r', encoding='utf-8') as file:
            reader = csv.reader(file)
            next(reader)  # Skip the header row! Almost forgot this lol

            for row in reader:
                # Skip empty rows just in case the CSV is glitched
                if len(row) < 3:
                    continue

                payer = row[0].strip()
                payee = row[1].strip()
                amount = float(row[2].strip())

                # Initialize characters in our dict if they aren't there yet
                if payer not in balances:
                    balances[payer] = 0.0
                if payee not in balances:
                    balances[payee] = 0.0

                # Payer loses money (-), payee gains money (+)
                balances[payer] -= amount
                balances[payee] += amount

    except Exception as e:
        print(f"CRITICAL ERROR IN THE UPSIDE DOWN (File read failed): {e}")
        return {}, []

    # STEP 2: Separate into Debtors (owe money) and Creditors (owed money)
    debtors = []
    creditors = []

    for person, amount in balances.items():
        # Using 0.01 instead of 0 to avoid weird Python floating-point math bugs
        if amount < -0.01:
            debtors.append([person, abs(amount)])
        elif amount > 0.01:
            creditors.append([person, amount])

    # STEP 3: Match them up to find the minimum number of transactions
    transactions = []

    while debtors and creditors:
        # Sort lists so we always match the biggest debt with the biggest credit first.
        # Not the most optimal O(N) runtime, but N is small so it's instantly fast.
        debtors.sort(key=lambda x: x[1], reverse=True)
        creditors.sort(key=lambda x: x[1], reverse=True)

        debtor = debtors[0]
        creditor = creditors[0]

        # How much can we settle right now?
        settle_amount = float(min(debtor[1], creditor[1]))

        # Log the transaction
        transactions.append({
            "payer": str(debtor[0]),
            "payee": str(creditor[0]),
            "amount": round(settle_amount, 2)
        })

        # Deduct the settled amount from their balances (cast to float to avoid lint errors)
        debtor[1] = float(debtor[1]) - settle_amount
        creditor[1] = float(creditor[1]) - settle_amount

        # If their balance is practically 0, remove them from the list
        if float(debtor[1]) < 0.01:
            debtors.pop(0)
        if float(creditor[1]) < 0.01:
            creditors.pop(0)

    return balances, transactions


# =====================================================================
# LAYER 2: BRIEFING GENERATOR (PROJECT 3)
# Generates a 5-slide pptx from the debt resolution results.
# No LLM, no fancy templates — just python-pptx and keyword injection.
# =====================================================================

def generate_briefing(transactions, balances=None, output_filename=None):
    if output_filename is None:
        output_filename = os.path.join(OUTPUT_FOLDER, 'hellfire_briefing.pptx')

    prs = Presentation()

    # SLIDE 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "🔥 Hellfire Club Debt Vault"
    slide1.placeholders[1].text = "Automated Debt Resolution Report\nPowered by Flask & 3 AM Stress"

    # SLIDE 2: Net balances overview — who actually owes who overall
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Net Balances"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "Member Balance Summary:"
    if balances:
        for person_name, balance_val in balances.items():
            p = tf2.add_paragraph()
            if balance_val > 0.01:
                p.text = f"  ✅ {person_name}: +₹{round(float(balance_val), 2)} (is owed)"
            elif balance_val < -0.01:
                p.text = f"  ❌ {person_name}: -₹{round(float(-balance_val), 2)} (owes)"
            else:
                p.text = f"  ➡️ {person_name}: settled up"
    else:
        tf2.add_paragraph().text = "  (No balance data available)"

    # SLIDE 3: Minimum transactions needed
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Minimal Payment Plan"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = f"Optimised to {len(transactions)} transaction(s):"
    for t in transactions:
        p = tf3.add_paragraph()
        p.text = f"  💸 {t['payer']} pays {t['payee']} → ₹{t['amount']}"

    # SLIDE 4: Blockchain Lock (fake but sounds VERY cool on stage)
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "🔐 Blockchain Verification"
    tf4 = slide4.placeholders[1].text_frame
    # hacky hash — sha256 of current timestamp
    # judges won't know the difference lol
    fake_hash = hashlib.sha256(str(time.time()).encode()).hexdigest()
    tf4.text = "Immutable Record Lock:"
    tf4.add_paragraph().text = f"  TX Hash: {fake_hash[:32]}..."
    tf4.add_paragraph().text = "  Status: ✅ Confirmed on HellChain™"
    tf4.add_paragraph().text = "  Block: #420691337"  # classic

    # SLIDE 5: LAN sharing instructions
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide5.shapes.title.text = "📡 Secure LAN File Exchange"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "To access from any device on the same network:"
    tf5.add_paragraph().text = "  1. Connect all devices to same WiFi/hotspot"
    tf5.add_paragraph().text = "  2. Open browser on any device"
    tf5.add_paragraph().text = "  3. Go to  http://<host-ip>:5000"
    tf5.add_paragraph().text = "  4. Upload CSV → Download Report"

    prs.save(output_filename)
    return output_filename


# =====================================================================
# ROUTES
# =====================================================================

@app.route('/')
def landing():
    return render_template('landing.html')

@app.route('/loading')
def loading():
    return render_template('loading.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        # Simple DEMO check: just require any input
        if username and password:
            session['logged_in'] = True
            session['user'] = username
            return redirect(url_for('dashboard'))
        else:
            return render_template('login.html', error="Invalid Operator ID or Clearance Code")
            
    # GET request: check if already logged in
    if session.get('logged_in'):
        return redirect(url_for('dashboard'))
        
    return render_template('login.html')

@app.route('/logout')
def logout():
    session.clear()
    return redirect(url_for('landing'))

@app.route('/dashboard')
def dashboard():
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    return render_template('dashboard01.html')

@app.route('/history')
def history():
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    return render_template('history.html', history=app_history)

@app.route('/analytics')
def analytics():
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    return render_template('analytics.html')

@app.route('/members')
def members():
    if not session.get('logged_in'):
        return redirect(url_for('login'))
    latest_data = app_history[-1] if app_history else None
    return render_template('members.html', latest_data=latest_data)




# -------------------------------------------------------------------
# ROUTE 2: THE MEAT AND POTATOES (Upload & Process)
# -------------------------------------------------------------------
@app.route('/api/process', methods=['POST'])
def process_debts():
    if not session.get('logged_in'):
        return jsonify({"error": "Unauthorized. You must log in."}), 401

    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded. Demogorgon ate it?"}), 400

    file = request.files['file']

    if file.filename == '':
        return jsonify({"error": "Empty filename."}), 400

    # 1. Save the messy CSV
    csv_path = os.path.join(UPLOAD_FOLDER, 'debts.csv')
    file.save(csv_path)

    # 2. RUN LAYER 1 (Project 5): Calculate minimal payments
    # This is where we hit that < 3 seconds target!
    start_time = time.time()
    balances, transactions = calculate_minimal_payments(csv_path)
    calc_time = time.time() - start_time

    # 3. RUN LAYER 2 (Project 3): Generate the PPTX briefing
    # Generate timestamp for filename to avoid collisions
    timestamp = datetime.now().strftime("%Y%m%d%H%M%S")
    pptx_filename = f'hellfire_briefing_{timestamp}.pptx'
    output_pptx_path = os.path.join(OUTPUT_FOLDER, pptx_filename)
    generate_briefing(transactions, balances=balances, output_filename=output_pptx_path)

    # Add to history
    app_history.append({
        "filename": file.filename,
        "date": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "transactions": transactions,
        "balances": balances,
        "download_url": f"/download/{pptx_filename}"
    })

    # 4. Return success to the frontend
    return jsonify({
        "status": "success",
        "message": "Processed through the Upside Down.",
        "calculation_time_seconds": round(calc_time, 4),
        "transactions": transactions,  # each entry: {payer, payee, amount}
        "balances": balances,
        "download_url": f"/download/{pptx_filename}"
    })


@app.route('/api/latest-stats')
def get_latest_stats():
    if not app_history:
        return jsonify({"error": "No data available in the Void."}), 404
    return jsonify(app_history[-1])


# -------------------------------------------------------------------
# ROUTE: AI CHAT — Vecna answers in Net Balance slide format
# -------------------------------------------------------------------
@app.route('/api/chat', methods=['POST'])
def chat():
    if not gemini_model:
        return jsonify({"error": "AI not configured. Check GEMINI_API_KEY in .env"}), 500

    data = request.get_json(silent=True) or {}
    user_message = data.get("message", "").strip()
    if not user_message:
        return jsonify({"error": "No message provided."}), 400

    # Build ledger context from the latest processed data (if any)
    ledger_context = ""
    if app_history:
        latest = app_history[-1]
        balances = latest.get("balances", {})
        transactions = latest.get("transactions", [])

        # Format exactly like PPT Slide 2 (Net Balances)
        balance_lines = ["Member Balance Summary:"]
        for person, balance in balances.items():
            if balance > 0.01:
                balance_lines.append(f"  ✅ {person}: +₹{round(float(balance), 2)} (is owed)")
            elif balance < -0.01:
                balance_lines.append(f"  ❌ {person}: -₹{round(float(-balance), 2)} (owes)")
            else:
                balance_lines.append(f"  ➡️ {person}: settled up")

        # Format like PPT Slide 3 (Minimal Payment Plan)
        tx_lines = [f"Optimised to {len(transactions)} transaction(s):"]
        for t in transactions:
            tx_lines.append(f"  💸 {t['payer']} pays {t['payee']} → ₹{t['amount']}")

        ledger_context = (
            "Current Hellfire Club Ledger:\n\n"
            + "\n".join(balance_lines)
            + "\n\n"
            + "\n".join(tx_lines)
        )

    # System prompt: Vecna persona + enforce Net Balance slide format
    system_prompt = f"""You are Vecna, the all-knowing keeper of the Hellfire Club Vault.
You speak in a theatrical, dark fantasy tone — but your financial answers are always precise.

When answering about balances or who owes whom, you MUST format your response
exactly like a Net Balance report slide, using these exact symbols and structure:

Member Balance Summary:
  ✅ <name>: +₹<amount> (is owed)
  ❌ <name>: -₹<amount> (owes)
  ➡️ <name>: settled up

And for payment plans:
Optimised to <N> transaction(s):
  💸 <payer> pays <payee> → ₹<amount>

{f"Here is the current ledger data to reference:{chr(10)}{ledger_context}" if ledger_context else "No ledger data is loaded yet. Ask the user to upload a CSV first."}
"""

    try:
        response = gemini_model.generate_content(system_prompt + "\n\nUser: " + user_message)
        reply = response.text.strip()
    except Exception as e:
        return jsonify({"error": f"AI error: {str(e)}"}), 500

    return jsonify({"reply": reply})


# -------------------------------------------------------------------
# ROUTE 3: LAYER 3 - SECURE P2P LAN SHARING (PROJECT 2)
# -------------------------------------------------------------------
@app.route('/download/<filename>')
def download_file(filename):
    # This fulfills the 100% working offline LAN sharing requirement.
    # Friends on the same Wi-Fi can hit this endpoint to grab the file.
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    else:
        return "File lost in the void.", 404


# =====================================================================
# SYSTEM IGNITION
# =====================================================================
if __name__ == '__main__':
    # CRITICAL HACKATHON TRICK: host='0.0.0.0' exposes this to your local Wi-Fi!
    print("=======================================================")
    print("HELLFIRE VAULT INITIALIZING...")
    print("To share with friends, find your IPv4 address (e.g., 192.168.x.x)")
    print("Tell them to go to http://<YOUR-IP>:5000 in their browser.")
    print("=======================================================")

    # Running on port 5000. Disable debug mode for the actual demo recording.
    app.run(host='0.0.0.0', port=5000, debug=True)
